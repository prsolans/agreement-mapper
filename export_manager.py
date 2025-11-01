"""
Export Manager: Handle PNG, PowerPoint, and Word exports of visualizations
"""
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from io import BytesIO
from pathlib import Path
import os
import math
from tempfile import NamedTemporaryFile


class ExportManager:
    """Handles exporting visualizations to PNG and PowerPoint"""

    @staticmethod
    def capture_html_as_png(
        html_content: str,
        width: int = 2400,
        height: int = 2400,
        scale: int = 2
    ) -> bytes:
        """
        Capture HTML content as high-resolution PNG

        Args:
            html_content: Complete HTML string with embedded JavaScript
            width: Viewport width in pixels
            height: Viewport height in pixels
            scale: Device scale factor (2 = retina/high DPI)

        Returns:
            PNG image as bytes
        """
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(
                viewport={'width': width, 'height': height},
                device_scale_factor=scale
            )

            # Load HTML content
            page.set_content(html_content)

            # Wait for network idle (important for D3.js and dynamic content)
            page.wait_for_load_state('networkidle')

            # Additional wait for animations and rendering
            page.wait_for_timeout(2000)  # 2 seconds

            # Capture screenshot as bytes
            screenshot_bytes = page.screenshot(full_page=True, type='png')

            browser.close()

            return screenshot_bytes

    @staticmethod
    def create_pptx_with_visualization(
        company_name: str,
        image_bytes: bytes,
        analysis_data: dict
    ) -> bytes:
        """
        Create PowerPoint presentation with visualization

        Args:
            company_name: Company name for title
            image_bytes: PNG image bytes
            analysis_data: Full analysis dictionary

        Returns:
            PowerPoint file as bytes
        """
        # Create new presentation
        prs = Presentation()

        # Slide 1: Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"{company_name} Agreement Landscape"

        meta = analysis_data.get('_meta', {})
        analysis_date = meta.get('analysis_date', 'N/A')
        subtitle.text = f"Analysis Date: {analysis_date}"

        # Slide 2: Visualization (full-screen image)
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Save image bytes to temporary file for adding to slide
        with NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_bytes)
            tmp_path = tmp_file.name

        # Add image to slide (centered, 9" wide)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)  # 10" - 1" margins

        slide.shapes.add_picture(tmp_path, left, top, width=width)

        # Clean up temp file
        os.unlink(tmp_path)

        # Slide 3: Summary metrics
        title_content_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = "Key Insights"

        # Add text box with key metrics
        content = slide.placeholders[1]
        tf = content.text_frame

        profile = analysis_data.get('company_profile', {})
        scale = profile.get('scale', {})
        business_units = analysis_data.get('business_units', [])
        opportunities = analysis_data.get('optimization_opportunities', [])

        tf.text = f"Revenue: {scale.get('annual_revenue', 'N/A')}"

        p = tf.add_paragraph()
        p.text = f"Employees: {scale.get('employees', 'N/A')}"

        p = tf.add_paragraph()
        p.text = f"Business Units: {len(business_units)}"

        p = tf.add_paragraph()
        p.text = f"Optimization Opportunities: {len(opportunities)}"

        portfolio = analysis_data.get('portfolio_summary', {})
        if portfolio:
            p = tf.add_paragraph()
            p.text = f"Total Annual Value: {portfolio.get('total_annual_value', 'N/A')}"

            p = tf.add_paragraph()
            p.text = f"Portfolio ROI: {portfolio.get('portfolio_roi', 'N/A')}"

        # Save to bytes
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        return pptx_bytes.getvalue()

    @staticmethod
    def create_pptx_native(
        company_name: str,
        analysis_data: dict
    ) -> bytes:
        """
        Create PowerPoint presentation with native shapes (no images)

        Args:
            company_name: Company name for title
            analysis_data: Full analysis dictionary

        Returns:
            PowerPoint file as bytes
        """
        # Create new presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Widescreen 16:9
        prs.slide_height = Inches(7.5)

        # Extract data
        profile = analysis_data.get('company_profile', {})
        scale = profile.get('scale', {})
        business_units = analysis_data.get('business_units', [])
        landscape = analysis_data.get('agreement_landscape_by_function', {})
        functions = landscape.get('functions', []) if isinstance(landscape, dict) else []
        opportunities = analysis_data.get('optimization_opportunities', [])
        agreement_matrix = analysis_data.get('agreement_matrix', {})
        matrix_types = agreement_matrix.get('agreement_types', [])
        meta = analysis_data.get('_meta', {})
        portfolio = analysis_data.get('portfolio_summary', {})
        priority_mappings = analysis_data.get('priority_mappings', [])

        # Standard DocuSign colors (blues and purples, dark enough for white text)
        def hex_to_rgb(hex_color):
            """Convert hex color to RGBColor"""
            hex_color = hex_color.lstrip('#')
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

        # DocuSign brand colors - blues and purples
        primary_color = hex_to_rgb('#E91E8C')      # Deep magenta/pink
        secondary_color = hex_to_rgb('#8B5CF6')    # Purple/violet
        accent_color = hex_to_rgb('#5B21B6')       # Deep blue-purple

        # ===== SLIDE 1: COMPANY HEADER & PROFILE =====
        blank_slide = prs.slide_layouts[6]
        slide1 = prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"{company_name} Agreement Landscape"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 121)

        # Subtitle
        subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        analysis_date = meta.get('analysis_date', 'N/A')
        subtitle_frame.text = f"Analysis Date: {analysis_date}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(14)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)

        # Key Metrics Grid
        metrics = [
            ("Industry", profile.get('industry', 'N/A')),
            ("Revenue", scale.get('annual_revenue', 'N/A')),
            ("Employees", str(scale.get('employees', 'N/A'))),
            ("Founded", str(profile.get('founded_year', 'N/A'))),
            ("Business Units", str(len(business_units))),
            ("Functions", str(len(functions))),
            ("Total Annual Value", portfolio.get('total_annual_value', 'N/A')),
            ("Portfolio ROI", portfolio.get('portfolio_roi', 'N/A'))
        ]

        y_pos = 2.2
        for i in range(0, len(metrics), 2):
            # Left metric
            if i < len(metrics):
                label_box = slide1.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.4))
                label_frame = label_box.text_frame
                label_frame.text = metrics[i][0]
                label_frame.paragraphs[0].font.size = Pt(12)
                label_frame.paragraphs[0].font.bold = True
                label_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

                value_box = slide1.shapes.add_textbox(Inches(3.0), Inches(y_pos), Inches(2.5), Inches(0.4))
                value_frame = value_box.text_frame
                value_frame.text = metrics[i][1]
                value_frame.paragraphs[0].font.size = Pt(12)
                value_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

            # Right metric
            if i + 1 < len(metrics):
                label_box = slide1.shapes.add_textbox(Inches(7.0), Inches(y_pos), Inches(2.5), Inches(0.4))
                label_frame = label_box.text_frame
                label_frame.text = metrics[i + 1][0]
                label_frame.paragraphs[0].font.size = Pt(12)
                label_frame.paragraphs[0].font.bold = True
                label_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

                value_box = slide1.shapes.add_textbox(Inches(9.5), Inches(y_pos), Inches(2.5), Inches(0.4))
                value_frame = value_box.text_frame
                value_frame.text = metrics[i + 1][1]
                value_frame.paragraphs[0].font.size = Pt(12)
                value_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

            y_pos += 0.6

        # ===== SLIDE 2: AGREEMENT LANDSCAPE BY FUNCTION (DAY FORCE STYLE) =====
        slide2 = prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"The organization's highly decentralized structure has led to an equally fragmented agreement landscape"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(18)
        title_para.font.bold = True

        # Position functions in scattered organic layout
        # Predefined positions for up to 8 functions (scattered arrangement) - adjusted for 13.33" width
        positions = [
            (Inches(1.5), Inches(1.5)),    # Top left
            (Inches(8.5), Inches(1.5)),    # Top right
            (Inches(0.8), Inches(3.6)),    # Middle left
            (Inches(6.0), Inches(3.5)),    # Center
            (Inches(10.5), Inches(3.6)),   # Middle right (box ends at ~13.3")
            (Inches(1.5), Inches(5.7)),    # Bottom left
            (Inches(8.5), Inches(5.7)),    # Bottom right
            (Inches(4.75), Inches(5.7)),   # Bottom center
        ]

        # Color rotation using brand colors
        box_colors = [primary_color, secondary_color, primary_color, secondary_color,
                      primary_color, secondary_color, primary_color, secondary_color]

        num_functions = min(len(functions), 8)  # Limit to 8 for layout
        if num_functions > 0:
            for idx, func in enumerate(functions[:num_functions]):
                pos_x, pos_y = positions[idx]

                # Box dimensions
                box_width = Inches(2.8)
                box_height = Inches(1.4)

                # Helper function to lighten a color
                def lighten_color(rgb_color, factor=0.4):
                    """Lighten an RGBColor by mixing with white"""
                    # RGBColor can be unpacked as a tuple (r, g, b)
                    r_orig, g_orig, b_orig = rgb_color
                    r = int(r_orig + (255 - r_orig) * factor)
                    g = int(g_orig + (255 - g_orig) * factor)
                    b = int(b_orig + (255 - b_orig) * factor)
                    return RGBColor(r, g, b)

                # LAYER 1: Light-colored background (full box)
                light_bg = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    pos_x,
                    pos_y,
                    box_width,
                    box_height
                )
                light_bg.fill.solid()
                light_bg.fill.fore_color.rgb = lighten_color(box_colors[idx])
                light_bg.line.width = Pt(0)  # No border

                # LAYER 2: Dark-colored section (right portion)
                dark_width = Inches(1.8)  # Right section width
                dark_offset = Inches(1.0)  # Start position from left
                dark_box = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    pos_x + dark_offset,
                    pos_y,
                    dark_width,
                    box_height
                )
                dark_box.fill.solid()
                dark_box.fill.fore_color.rgb = box_colors[idx]
                dark_box.line.width = Pt(0)  # No border

                # LAYER 3: White rounded rectangle (left section, on top)
                white_width = Inches(1.2)
                white_height = Inches(0.8)
                white_offset_x = Inches(0.15)
                white_offset_y = (box_height - white_height) / 2

                white_rect = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    pos_x + white_offset_x,
                    pos_y + white_offset_y,
                    white_width,
                    white_height
                )
                white_rect.fill.solid()
                white_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
                white_rect.line.color.rgb = RGBColor(230, 230, 230)
                white_rect.line.width = Pt(0.5)

                # Function name in white rounded rectangle (BLACK TEXT)
                oval_frame = white_rect.text_frame
                oval_frame.word_wrap = True
                oval_frame.margin_left = Pt(4)
                oval_frame.margin_right = Pt(4)
                oval_frame.margin_top = Pt(4)
                oval_frame.margin_bottom = Pt(4)
                oval_frame.vertical_anchor = MSO_ANCHOR.TOP

                func_name = func.get('function_name', 'Unknown')
                p = oval_frame.paragraphs[0]
                p.text = func_name
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)  # BLACK TEXT
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(2)

                # Business unit descriptor (blue text)
                business_units_list = func.get('business_units', [])
                if business_units_list:
                    descriptor = business_units_list[0] if isinstance(business_units_list[0], str) else "Operations"
                else:
                    descriptor = "Operations"

                p = oval_frame.add_paragraph()
                p.text = descriptor[:15]  # Limit length
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(65, 105, 225)  # Royal blue
                p.alignment = PP_ALIGN.CENTER

                # Right section content (on the colored background)
                # Top line: Commerce info and agreement count
                agreements = func.get('total_agreements', 'N/A')
                top_text = slide2.shapes.add_textbox(
                    pos_x + Inches(1.3),
                    pos_y + Inches(0.08),
                    Inches(1.35),
                    Inches(0.25)
                )
                top_frame = top_text.text_frame
                top_frame.word_wrap = False
                p = top_frame.paragraphs[0]
                p.text = f"{agreements}+ Agreements"
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True

                # System badges (dark pills on top right - 50% overlay)
                systems = func.get('systems_used', [])
                badge_width = Inches(0.5)
                # Position so 50% hangs off right edge: right_edge - (badge_width / 2)
                badge_x = pos_x + box_width - (badge_width / 2)
                badge_y = pos_y + Inches(0.08)

                for sys_idx, system in enumerate(systems[:3]):  # Max 3 systems
                    badge = slide2.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        badge_x,
                        badge_y + (sys_idx * Inches(0.22)),
                        badge_width,
                        Inches(0.18)
                    )
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = RGBColor(26, 58, 82)  # Dark navy
                    badge.line.color.rgb = RGBColor(26, 58, 82)

                    badge_frame = badge.text_frame
                    badge_frame.margin_left = Pt(2)
                    badge_frame.margin_right = Pt(2)
                    p = badge_frame.paragraphs[0]
                    p.text = system[:6]  # Abbreviate if needed
                    p.font.size = Pt(6)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Agreement details (middle section)
                agreement_types = func.get('agreement_types', [])
                if agreement_types:
                    details_text = []
                    for agr in agreement_types[:2]:  # Show top 2
                        agr_name = agr.get('type', '')
                        if agr_name:
                            # Extract key words
                            words = agr_name.split()
                            if len(words) > 2:
                                details_text.append(words[0] + " " + words[1])
                            else:
                                details_text.append(agr_name[:15])

                    details_box = slide2.shapes.add_textbox(
                        pos_x + Inches(1.3),
                        pos_y + Inches(0.4),
                        Inches(1.3),
                        Inches(0.5)
                    )
                    details_frame = details_box.text_frame
                    details_frame.word_wrap = True
                    p = details_frame.paragraphs[0]
                    p.text = ", ".join(details_text) if details_text else "Various Agreements"
                    p.font.size = Pt(7)
                    p.font.color.rgb = RGBColor(255, 255, 255)

                # Bottom: Complexity indicator with icon
                complexity = func.get('complexity', 3)
                complexity_text = "Complex, Negotiated" if complexity >= 4 else "Moderate Complexity"

                bottom_box = slide2.shapes.add_textbox(
                    pos_x + Inches(1.3),
                    pos_y + Inches(1.05),
                    Inches(1.3),
                    Inches(0.25)
                )
                bottom_frame = bottom_box.text_frame
                p = bottom_frame.paragraphs[0]
                # Font Awesome file icon
                p.text = f"\uf15b {complexity_text}"
                p.font.name = "Font Awesome 6 Free"
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(255, 255, 255)

        # ===== SLIDE 3: PRIORITIES MAPPED TO CAPABILITIES =====
        slide3 = prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Priorities Mapped to Capabilities"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Subtitle
        subtitle_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.33), Inches(0.3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Strategic Alignment"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)

        # Create priority mappings layout
        if priority_mappings:
            # Color scheme - use company brand colors
            colors = [primary_color, secondary_color, accent_color]

            # Layout settings
            start_y = Inches(1.4)
            row_height = Inches(1.8)

            # Left section (priority)
            priority_pill_width = Inches(2.5)
            priority_pill_height = Inches(0.4)
            priority_blob_width = Inches(5.0)
            priority_blob_height = Inches(1.0)
            priority_left = Inches(0.5)

            # Connector
            connector_x = priority_left + priority_blob_width + Inches(0.3)

            # Right section (capability)
            capability_pill_width = Inches(2.5)
            capability_pill_height = Inches(0.4)
            capability_blob_width = Inches(5.0)
            capability_blob_height = Inches(1.0)
            capability_left = connector_x + Inches(0.6)

            for idx, mapping in enumerate(priority_mappings[:3]):  # Show first 3 mappings
                # Get data
                priority_name = mapping.get('priority_name', f'Priority {idx + 1}')
                priority_description = mapping.get('priority_description', '')
                capability_name = mapping.get('capability_name', f'Capability {idx + 1}')
                capability_description = mapping.get('capability_description', '')

                # Calculate Y position for this row
                y_pos = start_y + (idx * row_height)

                # === LEFT SECTION: PRIORITY ===
                # 1. Priority name pill
                priority_pill = slide3.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    priority_left,
                    y_pos,
                    priority_pill_width,
                    priority_pill_height
                )
                priority_pill.fill.solid()
                priority_pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
                priority_pill.line.color.rgb = colors[idx]
                priority_pill.line.width = Pt(2)

                pill_frame = priority_pill.text_frame
                pill_frame.word_wrap = False
                p = pill_frame.paragraphs[0]
                p.text = priority_name
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER
                pill_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 2. Priority description blob (colored)
                priority_blob_y = y_pos + priority_pill_height + Inches(0.1)
                priority_blob = slide3.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    priority_left,
                    priority_blob_y,
                    priority_blob_width,
                    priority_blob_height
                )
                priority_blob.fill.solid()
                priority_blob.fill.fore_color.rgb = colors[idx]
                priority_blob.line.width = Pt(0)

                blob_frame = priority_blob.text_frame
                blob_frame.word_wrap = True
                blob_frame.margin_left = Pt(15)
                blob_frame.margin_right = Pt(15)
                blob_frame.margin_top = Pt(12)
                blob_frame.margin_bottom = Pt(12)
                blob_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = blob_frame.paragraphs[0]
                p.text = priority_description
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(255, 255, 255)

                # === CONNECTOR ===
                connector_y = y_pos + priority_pill_height + (priority_blob_height / 2)

                # Draw line
                connector_line = slide3.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    int(priority_left + priority_blob_width),
                    int(connector_y),
                    int(capability_left),
                    int(connector_y)
                )
                connector_line.line.color.rgb = colors[idx]
                connector_line.line.width = Pt(2)

                # Add circular node in the middle
                node_x = connector_x
                node_size = Inches(0.15)
                node = slide3.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    node_x - node_size / 2,
                    connector_y - node_size / 2,
                    node_size,
                    node_size
                )
                node.fill.solid()
                node.fill.fore_color.rgb = colors[idx]
                node.line.width = Pt(0)

                # === RIGHT SECTION: CAPABILITY ===
                # 3. Capability name pill
                capability_pill = slide3.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    capability_left,
                    y_pos,
                    capability_pill_width,
                    capability_pill_height
                )
                capability_pill.fill.solid()
                capability_pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
                capability_pill.line.color.rgb = colors[idx]
                capability_pill.line.width = Pt(2)

                cap_pill_frame = capability_pill.text_frame
                cap_pill_frame.word_wrap = False
                p = cap_pill_frame.paragraphs[0]
                p.text = capability_name
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER
                cap_pill_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 4. Capability description blob (white)
                capability_blob_y = y_pos + capability_pill_height + Inches(0.1)
                capability_blob = slide3.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    capability_left,
                    capability_blob_y,
                    capability_blob_width,
                    capability_blob_height
                )
                capability_blob.fill.solid()
                capability_blob.fill.fore_color.rgb = RGBColor(255, 255, 255)
                capability_blob.line.color.rgb = colors[idx]
                capability_blob.line.width = Pt(2)

                cap_blob_frame = capability_blob.text_frame
                cap_blob_frame.word_wrap = True
                cap_blob_frame.margin_left = Pt(15)
                cap_blob_frame.margin_right = Pt(15)
                cap_blob_frame.margin_top = Pt(12)
                cap_blob_frame.margin_bottom = Pt(12)
                cap_blob_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = cap_blob_frame.paragraphs[0]
                p.text = capability_description
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0, 0, 0)

        # Add executive quotes to slide notes
        strategic_priorities = analysis.get('strategic_priorities', [])
        if strategic_priorities:
            notes_content = ["EXECUTIVE QUOTES - Strategic Priorities\n" + "="*50 + "\n"]

            for priority in strategic_priorities:
                priority_name = priority.get('priority_name', '')
                executive_quotes = priority.get('executive_quotes', [])

                if executive_quotes:
                    notes_content.append(f"\n{priority_name}:")
                    for quote in executive_quotes:
                        exec_name = quote.get('executive', 'Unknown')
                        quote_text = quote.get('quote', '')
                        source = quote.get('source', '')
                        date = quote.get('date', '')
                        url = quote.get('url', '')

                        notes_content.append(f"\n  • \"{quote_text}\"")
                        notes_content.append(f"    - {exec_name}")
                        notes_content.append(f"    - {source}, {date}")
                        notes_content.append(f"    - {url}")

            if len(notes_content) > 1:  # If we have more than just the header
                notes_slide = slide3.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = "\n".join(notes_content)

        # ===== SLIDE 4: TOP 3 IDENTIFIED USE CASES =====
        slide4 = prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Top 3 Identified Use Cases"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Subtitle
        subtitle_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.33), Inches(0.3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Intelligent Agreement Management"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(16)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)

        # Create opportunities layout
        if opportunities:
            # Take up to 3 opportunities
            display_opps = opportunities[:3]

            # Color scheme - use company brand colors
            colors = [primary_color, secondary_color, accent_color]

            # Layout settings
            start_y = Inches(1.4)
            row_height = Inches(1.9)
            pill_width = Inches(2.5)
            pill_height = Inches(0.4)
            blob_left = Inches(0.5)
            blob_width = Inches(7.5)
            blob_height = Inches(1.5)
            metrics_left = Inches(8.5)
            metric_size = Inches(1.0)
            metric_spacing = Inches(1.3)

            for idx, opp in enumerate(display_opps):
                # Get data
                use_case_name = opp.get('use_case_name', opp.get('title', f"Use Case {idx + 1}"))
                capabilities = opp.get('capabilities', opp.get('description', 'N/A'))
                agreement_types = opp.get('agreement_types', [])
                risk_reduction = opp.get('risk_reduction', '')
                compliance = opp.get('compliance_benefits', '')
                metrics = opp.get('metrics', [])

                # Calculate Y position for this row
                y_pos = start_y + (idx * row_height)

                # 1. Use case name pill (top left)
                pill = slide4.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    blob_left,
                    y_pos,
                    pill_width,
                    pill_height
                )
                pill.fill.solid()
                pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
                pill.line.color.rgb = colors[idx]
                pill.line.width = Pt(2)

                pill_frame = pill.text_frame
                pill_frame.word_wrap = False
                p = pill_frame.paragraphs[0]
                p.text = use_case_name
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER
                pill_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 2. Content blob (colored rounded rectangle)
                blob_y = y_pos + pill_height + Inches(0.1)
                blob = slide4.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    blob_left,
                    blob_y,
                    blob_width,
                    blob_height
                )
                blob.fill.solid()
                blob.fill.fore_color.rgb = colors[idx]
                blob.line.width = Pt(0)

                blob_frame = blob.text_frame
                blob_frame.word_wrap = True
                blob_frame.margin_left = Pt(15)
                blob_frame.margin_right = Pt(15)
                blob_frame.margin_top = Pt(12)
                blob_frame.margin_bottom = Pt(12)

                # Capabilities
                p = blob_frame.paragraphs[0]
                p.text = "Capabilities: "
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                # Add capabilities text in same paragraph
                run = p.runs[0]
                run.text = "Capabilities: "
                new_run = p.add_run()
                new_run.text = capabilities
                new_run.font.size = Pt(8)
                new_run.font.bold = False
                new_run.font.color.rgb = RGBColor(255, 255, 255)
                p.space_after = Pt(6)

                # Agreement types
                if agreement_types:
                    p = blob_frame.add_paragraph()
                    p.text = f"Agreements: {', '.join(agreement_types)}"
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.space_after = Pt(6)

                # Risk & Compliance (if available)
                if risk_reduction or compliance:
                    risk_comp_text = []
                    if risk_reduction:
                        risk_comp_text.append(risk_reduction)
                    if compliance:
                        risk_comp_text.append(compliance)

                    p = blob_frame.add_paragraph()
                    p.text = " | ".join(risk_comp_text)
                    p.font.size = Pt(7)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.italic = True

                # 3. Metric circles (right side) - show first 4 metrics
                display_metrics = metrics[:4] if metrics else []
                for metric_idx, metric in enumerate(display_metrics):
                    # Calculate metric position
                    if metric_idx < 2:
                        # Top row
                        metric_x = metrics_left + (metric_idx * metric_spacing)
                        metric_y = y_pos + Inches(0.3)
                    else:
                        # Bottom row
                        metric_x = metrics_left + ((metric_idx - 2) * metric_spacing)
                        metric_y = y_pos + Inches(0.3) + metric_spacing

                    # Create circle
                    circle = slide4.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        metric_x,
                        metric_y,
                        metric_size,
                        metric_size
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    circle.line.color.rgb = colors[idx]
                    circle.line.width = Pt(3)

                    # Add metric text
                    circle_frame = circle.text_frame
                    circle_frame.word_wrap = True
                    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    circle_frame.margin_left = Pt(8)
                    circle_frame.margin_right = Pt(8)

                    # Metric value (large)
                    p = circle_frame.paragraphs[0]
                    p.text = metric.get('value', 'N/A')
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = colors[idx]
                    p.alignment = PP_ALIGN.CENTER
                    p.space_after = Pt(2)

                    # Metric label (small)
                    p = circle_frame.add_paragraph()
                    p.text = metric.get('label', '')
                    p.font.size = Pt(7)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.CENTER

        # Add executive alignment to slide notes
        opportunities = analysis.get('optimization_opportunities', [])
        if opportunities:
            notes_content = ["EXECUTIVE ALIGNMENT - Optimization Opportunities\n" + "="*50 + "\n"]

            for opp in opportunities[:3]:  # Top 3 opportunities
                opp_title = opp.get('title', 'Opportunity')
                exec_alignment = opp.get('executive_alignment', {})

                if exec_alignment:
                    priority_name = exec_alignment.get('priority_name', '')
                    exec_champion = exec_alignment.get('executive_champion', '')
                    alignment_statement = exec_alignment.get('alignment_statement', '')
                    supporting_quote = exec_alignment.get('supporting_quote', '')

                    notes_content.append(f"\n{opp_title}:")
                    if priority_name:
                        notes_content.append(f"  Supports Priority: {priority_name}")
                    if exec_champion:
                        notes_content.append(f"  Executive Champion: {exec_champion}")
                    if alignment_statement:
                        notes_content.append(f"  Alignment: {alignment_statement}")
                    if supporting_quote:
                        notes_content.append(f"  Supporting Quote: \"{supporting_quote}\"")

            if len(notes_content) > 1:  # If we have more than just the header
                notes_slide = slide4.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = "\n".join(notes_content)

        # ===== SLIDE 5: AGREEMENT LANDSCAPE MATRIX =====
        slide5 = prs.slides.add_slide(blank_slide)

        # Title
        title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = "Agreement Landscape Matrix"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # Create matrix visualization if data available
        if matrix_types:
            # Plot area dimensions - expanded for wider slide
            plot_left = Inches(1)
            plot_top = Inches(1.2)
            plot_width = Inches(9.5)  # Increased from 6.5" to 9.5"
            plot_height = Inches(5.5)
            plot_right = plot_left + plot_width
            plot_bottom = plot_top + plot_height

            # Draw background grid
            grid_color = RGBColor(220, 220, 220)

            # Vertical grid lines (volume axis)
            for i in range(11):  # 0-10
                x = plot_left + (i / 10 * plot_width)
                line = slide5.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    int(x),
                    int(plot_top),
                    int(x),
                    int(plot_bottom)
                )
                line.line.color.rgb = grid_color
                line.line.width = Pt(0.5)

            # Horizontal grid lines (complexity axis)
            for i in range(11):  # 0-10
                y = plot_bottom - (i / 10 * plot_height)
                line = slide5.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    int(plot_left),
                    int(y),
                    int(plot_right),
                    int(y)
                )
                line.line.color.rgb = grid_color
                line.line.width = Pt(0.5)

            # Axis labels
            # X-axis label
            x_label = slide5.shapes.add_textbox(plot_left, plot_bottom + Inches(0.1), plot_width, Inches(0.3))
            x_frame = x_label.text_frame
            x_frame.text = "Agreement Volume (Low → High)"
            x_frame.paragraphs[0].font.size = Pt(10)
            x_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Y-axis label (rotated effect with vertical text)
            y_label = slide5.shapes.add_textbox(Inches(0.3), plot_top, Inches(0.5), plot_height)
            y_frame = y_label.text_frame
            y_frame.text = "Complexity"
            y_frame.paragraphs[0].font.size = Pt(10)
            y_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            y_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Business unit color mapping
            bu_colors = {
                'Legal': RGBColor(70, 130, 180),      # Steel Blue
                'HR': RGBColor(34, 139, 34),           # Forest Green
                'Sales': RGBColor(220, 20, 60),        # Crimson
                'Procurement': RGBColor(255, 140, 0),  # Dark Orange
                'Operations': RGBColor(138, 43, 226),  # Blue Violet
                'IT': RGBColor(0, 191, 255),           # Deep Sky Blue
                'Finance': RGBColor(218, 165, 32),     # Goldenrod
                'Marketing': RGBColor(219, 112, 147)   # Pale Violet Red
            }
            default_color = RGBColor(128, 128, 128)  # Gray

            # Track used business units for legend
            used_business_units = set()

            # Plot agreement types as bubbles (limit to fit well)
            display_types = matrix_types[:18]  # Show up to 18 for readability

            for agr in display_types:
                volume = agr.get('volume', 5)
                complexity = agr.get('complexity', 5)
                agr_type = agr.get('type', 'Unknown')
                business_unit = agr.get('business_unit', 'Other')

                # Calculate position (normalize 1-10 to plot area)
                x_pos = plot_left + ((volume - 1) / 9 * plot_width)
                y_pos = plot_bottom - ((complexity - 1) / 9 * plot_height)

                # Uniform bubble size (volume is represented by X-axis position)
                bubble_size = Inches(0.85)

                # Get color for business unit
                bubble_color = bu_colors.get(business_unit, default_color)
                used_business_units.add(business_unit)

                # Create bubble
                bubble = slide5.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x_pos - bubble_size / 2,
                    y_pos - bubble_size / 2,
                    bubble_size,
                    bubble_size
                )
                bubble.fill.solid()
                bubble.fill.fore_color.rgb = bubble_color
                bubble.line.color.rgb = RGBColor(255, 255, 255)
                bubble.line.width = Pt(2)

                # Add label (full text, no truncation)
                label_frame = bubble.text_frame
                label_frame.word_wrap = True
                label_frame.margin_left = Pt(4)
                label_frame.margin_right = Pt(4)
                label_frame.margin_top = Pt(4)
                label_frame.margin_bottom = Pt(4)

                p = label_frame.paragraphs[0]
                p.text = agr_type  # Full text, no truncation
                p.font.size = Pt(7)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Add legend on the right - moved for wider slide
            legend_left = Inches(11.0)
            legend_top = Inches(1.5)
            legend_item_height = Inches(0.25)

            # Legend title
            legend_title = slide5.shapes.add_textbox(legend_left, Inches(1.2), Inches(2), Inches(0.3))
            legend_title_frame = legend_title.text_frame
            legend_title_frame.text = "Business Units"
            legend_title_frame.paragraphs[0].font.size = Pt(9)
            legend_title_frame.paragraphs[0].font.bold = True

            # Legend items
            y_offset = 0
            for bu in sorted(used_business_units):
                # Color box
                color_box = slide5.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    legend_left,
                    legend_top + Inches(y_offset),
                    Inches(0.2),
                    Inches(0.2)
                )
                color_box.fill.solid()
                color_box.fill.fore_color.rgb = bu_colors.get(bu, default_color)
                color_box.line.color.rgb = RGBColor(100, 100, 100)

                # Label
                label = slide5.shapes.add_textbox(
                    legend_left + Inches(0.25),
                    legend_top + Inches(y_offset),
                    Inches(1.7),
                    Inches(0.2)
                )
                label_frame = label.text_frame
                label_frame.text = bu
                label_frame.paragraphs[0].font.size = Pt(7)
                label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                y_offset += 0.28

        # Save to bytes
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        return pptx_bytes.getvalue()

    @staticmethod
    def create_docx_content(
        company_name: str,
        analysis_data: dict
    ) -> bytes:
        """
        Create Word document with slide content for copy/paste

        Args:
            company_name: Company name for title
            analysis_data: Full analysis dictionary

        Returns:
            Word document file as bytes
        """
        # Create new document
        doc = Document()

        # Extract data
        profile = analysis_data.get('company_profile', {})
        scale = profile.get('scale', {})
        business_units = analysis_data.get('business_units', [])
        landscape = analysis_data.get('agreement_landscape_by_function', {})
        functions = landscape.get('functions', []) if isinstance(landscape, dict) else []
        opportunities = analysis_data.get('optimization_opportunities', [])
        priority_mappings = analysis_data.get('priority_mappings', [])
        strategic_priorities = analysis_data.get('strategic_priorities', [])
        agreement_matrix = analysis_data.get('agreement_matrix', {})
        matrix_types = agreement_matrix.get('agreement_types', [])
        portfolio = analysis_data.get('portfolio_summary', {})

        # ===== SLIDE 1: COMPANY HEADER & PROFILE =====
        doc.add_heading(f'SLIDE 1: {company_name} Agreement Landscape', level=1)
        doc.add_paragraph(f'Company: {company_name}')
        doc.add_paragraph(f'Industry: {profile.get("industry", "N/A")}')
        doc.add_paragraph()

        # Metrics table
        table = doc.add_table(rows=6, cols=2)
        table.style = 'Light Grid Accent 1'

        # Header
        table.cell(0, 0).text = 'Metric'
        table.cell(0, 1).text = 'Value'

        # Data
        table.cell(1, 0).text = 'Annual Revenue'
        table.cell(1, 1).text = scale.get('annual_revenue', 'N/A')

        table.cell(2, 0).text = 'Employees'
        table.cell(2, 1).text = str(scale.get('employees', 'N/A'))

        table.cell(3, 0).text = 'Locations'
        table.cell(3, 1).text = str(scale.get('locations', 'N/A'))

        table.cell(4, 0).text = 'Countries'
        table.cell(4, 1).text = str(scale.get('countries', 'N/A'))

        table.cell(5, 0).text = 'Business Units'
        table.cell(5, 1).text = str(len(business_units))

        # Page break
        doc.add_page_break()

        # ===== SLIDE 2: AGREEMENT LANDSCAPE BY FUNCTION =====
        doc.add_heading('SLIDE 2: Agreement Landscape by Function', level=1)
        doc.add_paragraph('Business functions and their agreement management')
        doc.add_paragraph()

        if functions:
            # Functions table
            table = doc.add_table(rows=len(functions) + 1, cols=4)
            table.style = 'Light Grid Accent 1'

            # Header
            table.cell(0, 0).text = 'Function'
            table.cell(0, 1).text = 'Complexity'
            table.cell(0, 2).text = 'Total Agreements'
            table.cell(0, 3).text = 'Systems'

            # Data
            for idx, func in enumerate(functions, start=1):
                # Function name
                func_name = func.get('function_name', 'N/A')
                table.cell(idx, 0).text = func_name

                # Complexity (numeric, convert to text description)
                complexity = func.get('complexity', 3)
                complexity_text = "Complex, Negotiated" if complexity >= 4 else "Moderate Complexity"
                table.cell(idx, 1).text = complexity_text

                # Total agreements (volume)
                total_agreements = func.get('total_agreements', 'N/A')
                table.cell(idx, 2).text = str(total_agreements) if total_agreements != 'N/A' else 'N/A'

                # Systems used
                systems = func.get('systems_used', [])
                table.cell(idx, 3).text = ', '.join(systems) if systems else 'N/A'

        # Page break
        doc.add_page_break()

        # ===== SLIDE 3: PRIORITIES MAPPED TO CAPABILITIES =====
        doc.add_heading('SLIDE 3: Priorities Mapped to Capabilities', level=1)
        doc.add_paragraph('Strategic Alignment')
        doc.add_paragraph()

        if priority_mappings:
            for idx, mapping in enumerate(priority_mappings[:3], start=1):
                priority_name = mapping.get('priority_name', f'Priority {idx}')
                priority_description = mapping.get('priority_description', '')
                priority_id = mapping.get('priority_id', '')
                capability_name = mapping.get('capability_name', f'Capability {idx}')
                capability_description = mapping.get('capability_description', '')

                # Priority → Capability heading
                doc.add_heading(f'{idx}. {priority_name} → {capability_name}', level=2)

                # Priority description
                p = doc.add_paragraph()
                p.add_run('Priority: ').bold = True
                p.add_run(priority_description)

                # Capability description
                p = doc.add_paragraph()
                p.add_run('Capability: ').bold = True
                p.add_run(capability_description)

                doc.add_paragraph()

                # Find the full strategic priority data for notes
                priority_data = None
                if strategic_priorities:
                    for priority in strategic_priorities:
                        if priority.get('priority_id') == priority_id or priority.get('priority_name') == priority_name:
                            priority_data = priority
                            break

                # Add NOTES section with executive details
                if priority_data:
                    doc.add_heading('NOTES (Copy to Slide Notes):', level=3)

                    # Executive Owner
                    executive_owner = priority_data.get('executive_owner', 'N/A')
                    if executive_owner and executive_owner != 'N/A':
                        p = doc.add_paragraph()
                        p.add_run('Executive Owner: ').bold = True
                        p.add_run(executive_owner)

                    # Executive Responsibility
                    executive_responsibility = priority_data.get('executive_responsibility', 'N/A')
                    if executive_responsibility and executive_responsibility != 'N/A':
                        p = doc.add_paragraph()
                        p.add_run('Executive Responsibility: ').bold = True
                        p.add_run(executive_responsibility)

                    # Executive Quotes
                    executive_quotes = priority_data.get('executive_quotes', [])
                    if executive_quotes:
                        p = doc.add_paragraph()
                        p.add_run('Executive Quotes:').bold = True

                        for quote in executive_quotes:
                            exec_name = quote.get('executive', 'Unknown')
                            quote_text = quote.get('quote', '')
                            source = quote.get('source', '')
                            date = quote.get('date', '')
                            url = quote.get('url', '')

                            # Quote text
                            p = doc.add_paragraph(style='List Bullet')
                            p.add_run(f'"{quote_text}"').italic = True

                            # Attribution
                            p = doc.add_paragraph(style='List Bullet 2')
                            p.add_run(f'— {exec_name}, {source}, {date}')

                            # URL
                            if url:
                                p = doc.add_paragraph(style='List Bullet 2')
                                p.add_run(f'Source: {url}')

                    # Business Impact
                    business_impact = priority_data.get('business_impact', 'N/A')
                    if business_impact and business_impact != 'N/A':
                        p = doc.add_paragraph()
                        p.add_run('Business Impact: ').bold = True
                        p.add_run(business_impact)

                    # Related Initiatives
                    related_initiatives = priority_data.get('related_initiatives', [])
                    if related_initiatives:
                        p = doc.add_paragraph()
                        p.add_run('Related Initiatives: ').bold = True
                        if isinstance(related_initiatives, list):
                            p.add_run(', '.join(related_initiatives))
                        else:
                            p.add_run(str(related_initiatives))

                    # Sources
                    sources = priority_data.get('sources', [])
                    if sources:
                        p = doc.add_paragraph()
                        p.add_run('Sources: ').bold = True
                        if isinstance(sources, list):
                            p.add_run('; '.join(sources))
                        else:
                            p.add_run(str(sources))

                    doc.add_paragraph()

        # Page break
        doc.add_page_break()

        # ===== SLIDE 4: TOP 3 IDENTIFIED USE CASES =====
        doc.add_heading('SLIDE 4: Top 3 Identified Use Cases', level=1)
        doc.add_paragraph('Intelligent Agreement Management')
        doc.add_paragraph()

        if opportunities:
            for idx, opp in enumerate(opportunities[:3], start=1):
                use_case_name = opp.get('use_case_name', opp.get('title', f'Use Case {idx}'))
                capabilities = opp.get('capabilities', opp.get('description', 'N/A'))
                agreement_types = opp.get('agreement_types', [])
                risk_reduction = opp.get('risk_reduction', '')
                compliance = opp.get('compliance_benefits', '')
                metrics = opp.get('metrics', [])

                # Use case heading
                doc.add_heading(f'{idx}. {use_case_name}', level=2)

                # Capabilities
                p = doc.add_paragraph()
                p.add_run('Capabilities: ').bold = True
                p.add_run(capabilities)

                # Agreement types
                if agreement_types:
                    p = doc.add_paragraph()
                    p.add_run('Agreements: ').bold = True
                    p.add_run(', '.join(agreement_types))

                # Risk & Compliance
                if risk_reduction:
                    p = doc.add_paragraph()
                    p.add_run('Risk Reduction: ').bold = True
                    p.add_run(risk_reduction)

                if compliance:
                    p = doc.add_paragraph()
                    p.add_run('Compliance Benefits: ').bold = True
                    p.add_run(compliance)

                # Metrics
                if metrics:
                    doc.add_paragraph('Metrics:', style='Heading 3')
                    # Create 2x2 table for metrics
                    num_metrics = min(len(metrics), 4)
                    rows = 2 if num_metrics > 2 else 1
                    table = doc.add_table(rows=rows, cols=2)
                    table.style = 'Light Grid Accent 1'

                    for metric_idx, metric in enumerate(metrics[:4]):
                        row = metric_idx // 2
                        col = metric_idx % 2
                        cell = table.cell(row, col)
                        cell.text = f"{metric.get('value', 'N/A')}\n{metric.get('label', '')}"

                # Executive Alignment
                exec_alignment = opp.get('executive_alignment', {})
                if exec_alignment:
                    doc.add_paragraph('Executive Alignment:', style='Heading 3')

                    priority_name = exec_alignment.get('priority_name', '')
                    exec_champion = exec_alignment.get('executive_champion', '')
                    alignment_statement = exec_alignment.get('alignment_statement', '')
                    supporting_quote = exec_alignment.get('supporting_quote', '')

                    if priority_name:
                        p = doc.add_paragraph()
                        p.add_run('Supports Priority: ').bold = True
                        p.add_run(priority_name)

                    if exec_champion:
                        p = doc.add_paragraph()
                        p.add_run('Executive Champion: ').bold = True
                        p.add_run(exec_champion)

                    if alignment_statement:
                        p = doc.add_paragraph()
                        p.add_run('Alignment: ').bold = True
                        p.add_run(alignment_statement)

                    if supporting_quote:
                        p = doc.add_paragraph(style='List Bullet')
                        p.add_run(f'Supporting Quote: ').bold = True
                        p.add_run(f'"{supporting_quote}"').italic = True

                doc.add_paragraph()

        # Page break
        doc.add_page_break()

        # ===== SLIDE 5: AGREEMENT LANDSCAPE MATRIX =====
        doc.add_heading('SLIDE 5: Agreement Landscape Matrix', level=1)
        doc.add_paragraph('Agreement types by business unit, volume, and complexity')
        doc.add_paragraph()

        if matrix_types:
            # Matrix table
            display_types = matrix_types[:18]  # Show first 18
            table = doc.add_table(rows=len(display_types) + 1, cols=4)
            table.style = 'Light Grid Accent 1'

            # Header
            table.cell(0, 0).text = 'Agreement Type'
            table.cell(0, 1).text = 'Business Unit'
            table.cell(0, 2).text = 'Volume (1-10)'
            table.cell(0, 3).text = 'Complexity (1-10)'

            # Data
            for idx, agr in enumerate(display_types, start=1):
                table.cell(idx, 0).text = agr.get('type', 'Unknown')
                table.cell(idx, 1).text = agr.get('business_unit', 'Other')
                table.cell(idx, 2).text = str(agr.get('volume', 5))
                table.cell(idx, 3).text = str(agr.get('complexity', 5))

        # Save to bytes
        docx_bytes = BytesIO()
        doc.save(docx_bytes)
        docx_bytes.seek(0)

        return docx_bytes.getvalue()
